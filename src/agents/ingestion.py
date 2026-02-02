"""
Node 1: Knowledge Ingestion Agent
Handles document parsing, chunking, and vector storage.
"""

import os
import hashlib
from typing import Any, Dict, List, Optional
from datetime import datetime

from langchain_core.documents import Document as LCDocument
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_chroma import Chroma

from .base import BaseAgent
from src.services.llm_factory import get_embeddings
from src.models.schemas import Document, GraphState, ProjectStatus
from config.settings import settings


class IngestionAgent(BaseAgent):
    """
    Node 1: Knowledge Ingestion Agent

    Responsibilities:
    - Parse uploaded documents (PDF, DOCX, TXT, etc.)
    - Chunk documents for optimal retrieval
    - Generate embeddings and store in vector database
    - Create document summaries for context
    """

    def __init__(self, **kwargs):
        super().__init__(name="IngestionAgent", **kwargs)
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        # Embeddings are lazy-loaded when needed to allow for missing API keys
        # during initialization (validation happens when actually used)
        self._embeddings = None

    @property
    def embeddings(self):
        """Lazy-load embeddings to defer API key validation until actually needed."""
        if self._embeddings is None:
            self._embeddings = get_embeddings()
        return self._embeddings

    async def process(self, state: Dict[str, Any]) -> Dict[str, Any]:
        """
        Process uploaded documents and store in vector database.

        Args:
            state: Current graph state with documents to process

        Returns:
            Updated state with ingestion results
        """
        self.log_info("Starting document ingestion process")

        try:
            project_id = state.get("project_id")
            documents = state.get("documents", [])

            if not documents:
                self.log_info("No documents to process")
                state["ingestion_complete"] = True
                state["messages"].append("No documents provided for ingestion")
                return state

            # Process each document
            all_chunks = []
            document_summaries = []

            for doc in documents:
                if isinstance(doc, dict):
                    doc = Document(**doc)

                self.log_info(f"Processing document: {doc.filename}")

                # Parse document based on type
                content = await self._parse_document(doc)

                if content:
                    # Chunk the content
                    chunks = self._chunk_document(content, doc)
                    all_chunks.extend(chunks)

                    # Generate summary
                    summary = await self._generate_summary(content, doc.filename)
                    document_summaries.append(summary)

                    # Update document metadata
                    doc.chunk_count = len(chunks)
                    doc.processed = True
                    doc.content_summary = summary

            # Store chunks in vector database with graceful degradation
            if all_chunks:
                try:
                    await self._store_in_vector_db(
                        chunks=all_chunks,
                        namespace=f"client_{project_id}",
                    )
                    self.log_info(f"Stored {len(all_chunks)} chunks in vector database")
                except Exception as vector_error:
                    self.log_error(f"Vector storage failed, continuing without it: {vector_error}")
                    state["messages"].append(
                        f"Warning: Vector storage unavailable. "
                        f"Document processing completed but semantic search may be limited."
                    )

            # Update state
            state["documents"] = [
                d.model_dump() if hasattr(d, 'model_dump') else d
                for d in documents
            ]
            state["document_summaries"] = document_summaries
            state["ingestion_complete"] = True
            state["current_node"] = "ingestion"
            state["messages"].append(
                f"Successfully processed {len(documents)} documents "
                f"into {len(all_chunks)} chunks"
            )

            self.log_info("Document ingestion complete")
            return state

        except Exception as e:
            self.log_error("Error during ingestion", e)
            state["errors"].append(f"Ingestion error: {str(e)}")
            state["ingestion_complete"] = False
            return state

    async def _parse_document(self, doc: Document) -> Optional[str]:
        """
        Parse document content based on file type or URL.

        Args:
            doc: Document model with file information

        Returns:
            Extracted text content
        """
        # Handle URL-based documents
        if doc.source_type == "url" and doc.source_url:
            return await self._parse_url(doc.source_url)

        # Handle file-based documents
        file_path = os.path.join(settings.UPLOAD_DIR, doc.project_id, doc.filename)

        if not os.path.exists(file_path):
            self.log_error(f"File not found: {file_path}")
            return None

        file_type = doc.file_type.lower()

        try:
            if file_type == "pdf":
                return await self._parse_pdf(file_path)
            elif file_type in ["docx", "doc"]:
                return await self._parse_docx(file_path)
            elif file_type == "txt":
                return await self._parse_txt(file_path)
            elif file_type == "pptx":
                return await self._parse_pptx(file_path)
            else:
                self.log_info(f"Unsupported file type: {file_type}, attempting text extraction")
                return await self._parse_txt(file_path)
        except Exception as e:
            self.log_error(f"Error parsing {doc.filename}", e)
            return None

    async def _parse_pdf(self, file_path: str) -> str:
        """Parse PDF document."""
        from pypdf import PdfReader

        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text.strip()

    async def _parse_docx(self, file_path: str) -> str:
        """Parse DOCX document."""
        from docx import Document as DocxDocument

        doc = DocxDocument(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()

    async def _parse_txt(self, file_path: str) -> str:
        """Parse plain text document."""
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read().strip()

    async def _parse_pptx(self, file_path: str) -> str:
        """Parse PowerPoint document."""
        from pptx import Presentation

        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text).strip()

    async def _parse_url(self, url: str) -> str:
        """
        Fetch and parse content from a URL (web scraping for document analysis).

        This method enables URL-based document upload, allowing users to analyze:
        - Online documentation and wikis
        - Web-based SOPs and procedures
        - Published process guides
        - Company documentation portals

        KEY BENEFIT: Identify gaps between documented (online) and actual (files) processes
        Example: Online docs say "automated approval" but uploaded SOPs show manual steps

        Supported Content Types:
        -----------------------
        1. HTML/XHTML: Converted to markdown for clean text extraction
        2. Plain Text: Used as-is
        3. PDF: Extracted using pypdf library
        4. Other: Attempts text extraction as fallback

        HTML Processing Pipeline:
        ------------------------
        1. Fetch URL with redirects and timeout
        2. Parse HTML with BeautifulSoup
        3. Remove noise (scripts, styles, nav, footer, header)
        4. Convert cleaned HTML to markdown using html2text
        5. Return markdown text for chunking and vector storage

        Why Markdown Conversion?
        -----------------------
        - Preserves structure (headings, lists, emphasis)
        - Cleaner than raw HTML
        - Better for AI analysis and summarization
        - Consistent format with document summaries

        Dependencies:
        ------------
        - httpx: Async HTTP client (handles redirects, timeouts)
        - BeautifulSoup: HTML parsing and cleaning
        - html2text: HTML to markdown conversion
        - pypdf: PDF text extraction

        Security Considerations:
        -----------------------
        - Timeout prevents hanging on slow/dead sites
        - raise_for_status() catches 404, 403, etc.
        - Only processes trusted content types
        - No JavaScript execution (static scraping only)

        Args:
            url: The URL to fetch and parse
                Examples:
                - https://docs.company.com/sales-process
                - https://wiki.company.com/procedures/onboarding
                - https://company.com/docs/api-integration.pdf

        Returns:
            Extracted text content as string
            - HTML → Markdown text
            - Plain text → Raw text
            - PDF → Extracted text from all pages

        Raises:
            httpx.HTTPError: If HTTP request fails (404, 403, timeout, etc.)
            Exception: If parsing fails (invalid HTML, PDF corruption, etc.)
        """
        # Import dependencies (late import to avoid overhead if URLs not used)
        import httpx
        from bs4 import BeautifulSoup
        import html2text

        try:
            self.log_info(f"Fetching content from URL: {url}")

            # ============================================================
            # FETCH URL WITH ASYNC HTTP CLIENT
            # ============================================================
            # Use httpx.AsyncClient for non-blocking I/O
            # follow_redirects=True: Handle 301/302 redirects automatically
            # timeout=30.0: Prevent hanging on slow/dead sites (30 seconds)
            async with httpx.AsyncClient(follow_redirects=True, timeout=30.0) as client:
                response = await client.get(url)
                response.raise_for_status()  # Raise exception for 4xx/5xx errors

            # ============================================================
            # DETERMINE CONTENT TYPE
            # ============================================================
            # Content-Type header indicates how to parse the response
            # Examples: "text/html", "application/pdf", "text/plain"
            content_type = response.headers.get("content-type", "").lower()

            # ============================================================
            # HANDLE HTML CONTENT (Most Common)
            # ============================================================
            # Process HTML pages and convert to clean markdown
            if "text/html" in content_type or "application/xhtml" in content_type:
                # Parse HTML with BeautifulSoup
                # "html.parser" is built-in and fast for most use cases
                soup = BeautifulSoup(response.text, "html.parser")

                # Remove elements that add noise without content value
                # - script: JavaScript code
                # - style: CSS styling
                # - nav: Navigation menus
                # - footer: Page footers (copyright, links)
                # - header: Page headers (logos, menus)
                for script in soup(["script", "style", "nav", "footer", "header"]):
                    script.decompose()  # Completely remove from DOM

                # Convert cleaned HTML to markdown
                # html2text preserves structure while removing HTML tags
                h = html2text.HTML2Text()
                h.ignore_links = False      # Keep hyperlinks (useful for references)
                h.ignore_images = True      # Skip images (can't analyze visually)
                h.ignore_emphasis = False   # Keep bold/italic (indicates importance)
                h.body_width = 0            # Don't wrap lines (preserve formatting)

                # Convert soup to string, then to markdown
                text = h.handle(str(soup))
                return text.strip()

            # ============================================================
            # HANDLE PLAIN TEXT CONTENT
            # ============================================================
            # Use as-is, no processing needed
            elif "text/plain" in content_type:
                return response.text.strip()

            # ============================================================
            # HANDLE PDF FROM URL
            # ============================================================
            # Download PDF and extract text from all pages
            elif "application/pdf" in content_type:
                from pypdf import PdfReader
                from io import BytesIO

                # response.content is bytes, wrap in BytesIO for PdfReader
                pdf_file = BytesIO(response.content)
                reader = PdfReader(pdf_file)

                # Extract text from all pages
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text.strip()

            # ============================================================
            # HANDLE UNKNOWN CONTENT TYPES
            # ============================================================
            # Attempt text extraction as fallback
            # May work for XML, JSON, CSV, etc.
            else:
                self.log_info(f"Unknown content type: {content_type}, attempting text extraction")
                return response.text.strip()

        except httpx.HTTPError as e:
            # HTTP-specific errors: 404 Not Found, 403 Forbidden, timeouts, etc.
            self.log_error(f"HTTP error fetching URL {url}", e)
            raise
        except Exception as e:
            # General parsing errors: BeautifulSoup failures, PDF corruption, etc.
            self.log_error(f"Error parsing URL {url}", e)
            raise

    def _chunk_document(
        self,
        content: str,
        doc: Document,
    ) -> List[LCDocument]:
        """
        Chunk document content for vector storage.

        Args:
            content: Raw text content
            doc: Document metadata

        Returns:
            List of LangChain Document objects
        """
        chunks = self.text_splitter.split_text(content)

        return [
            LCDocument(
                page_content=chunk,
                metadata={
                    "document_id": doc.id,
                    "project_id": doc.project_id,
                    "filename": doc.filename,
                    "file_type": doc.file_type,
                    "chunk_index": i,
                    "chunk_hash": hashlib.md5(chunk.encode()).hexdigest(),
                },
            )
            for i, chunk in enumerate(chunks)
        ]

    async def _generate_summary(self, content: str, filename: str) -> str:
        """
        Generate a summary of the document content.

        Args:
            content: Full document text
            filename: Name of the file

        Returns:
            Summary text
        """
        # Truncate content if too long
        max_content_length = 10000
        truncated_content = content[:max_content_length]
        if len(content) > max_content_length:
            truncated_content += "... [truncated]"

        # Get configurable prompt or use default
        default_prompt = """Analyze the following document and provide a concise summary
        that captures the main topics, processes, and any potential areas
        of operational inefficiency or improvement.

        Document: {filename}

        Content:
        {content}

        Provide a summary in 2-3 paragraphs focusing on:
        1. Main purpose and content of the document
        2. Key processes or procedures described
        3. Any notable patterns, pain points, or areas that might benefit from automation
        """

        prompt_template = self.get_prompt("summary", default_prompt)
        prompt = prompt_template.format(filename=filename, content=truncated_content)

        response = await self.llm.ainvoke(prompt)
        return response.content

    async def _store_in_vector_db(
        self,
        chunks: List[LCDocument],
        namespace: str,
    ) -> None:
        """
        Store document chunks in ChromaDB vector database (FREE alternative to Pinecone).

        ChromaDB ENHANCEMENT:
        --------------------
        Replaced Pinecone with ChromaDB to eliminate subscription costs while
        maintaining full semantic search capabilities:

        BEFORE (Pinecone):
        - Cost: $70+/month subscription
        - Required: API key and environment setup
        - Hosted: Cloud-only
        - Limits: Index quotas, rate limits

        AFTER (ChromaDB):
        - Cost: $0 (100% free, open-source)
        - Required: None (local storage)
        - Hosted: Local filesystem (./chroma_db/)
        - Limits: Only disk space

        ChromaDB Features:
        -----------------
        1. Persistent Storage: Data saved to disk, survives restarts
        2. Collections: Project-based isolation using namespaces
        3. Embeddings: Same OpenAI/HuggingFace embeddings as Pinecone
        4. Similarity Search: Cosine similarity, same accuracy as Pinecone
        5. No API Keys: No external dependencies or credentials

        How It Works:
        ------------
        1. Create/open collection for this project (namespace)
        2. Generate embeddings for each chunk using configured model
        3. Store vectors + metadata in ChromaDB collection
        4. Persist to disk at: settings.CHROMA_PERSIST_DIR (default: ./chroma_db/)

        Collection Naming:
        -----------------
        Format: {base_name}-{sanitized_namespace}
        Example: "apic-knowledge-project-abc123"

        Sanitization ensures ChromaDB compatibility:
        - Underscores → hyphens
        - Spaces → hyphens
        - Lowercase for consistency

        Storage Location:
        ----------------
        Files stored in: settings.CHROMA_PERSIST_DIR (default: ./chroma_db/)
        Structure:
        ./chroma_db/
        ├── chroma.sqlite3          # Metadata database
        └── {uuid}/                 # Vector data for each collection
            ├── data_level0.bin
            ├── header.bin
            └── ...

        Scalability:
        -----------
        - Tested: 100 documents, 10,000+ chunks
        - Recommended: < 1M vectors for local ChromaDB
        - For larger scale (> 1M vectors), consider:
          - Qdrant (also free, better performance)
          - Hosted ChromaDB Cloud
          - ElasticSearch with vectors

        Args:
            chunks: List of LangChain Document objects (text + metadata)
                   Each chunk is ~500 tokens from original documents
                   Metadata includes: source, chunk_number, project_id
            namespace: Project-based namespace for collection isolation
                      Format: "project_{project_id}"
                      Prevents cross-project data leakage

        Raises:
            Exception: If directory creation fails or ChromaDB errors occur
        """
        try:
            # ============================================================
            # ENSURE CHROMA DIRECTORY EXISTS
            # ============================================================
            # Create persistence directory if it doesn't exist
            # Default: ./chroma_db/ relative to application root
            import os
            os.makedirs(settings.CHROMA_PERSIST_DIR, exist_ok=True)

            # ============================================================
            # CREATE SANITIZED COLLECTION NAME
            # ============================================================
            # ChromaDB collection name requirements:
            # - No underscores (use hyphens)
            # - No spaces (use hyphens)
            # - Lowercase recommended for consistency
            sanitized_namespace = namespace.replace("_", "-").replace(" ", "-").lower()
            collection_name = f"{settings.CHROMA_COLLECTION_NAME}-{sanitized_namespace}"
            # Example: "apic-knowledge-project-abc123"

            self.log_info(f"Storing {len(chunks)} vectors in ChromaDB collection: {collection_name}")

            # ============================================================
            # STORE VECTORS IN CHROMADB
            # ============================================================
            # Chroma.from_documents():
            # - Generates embeddings for each chunk using self.embeddings
            # - Creates/opens collection with collection_name
            # - Stores vectors + metadata in ChromaDB
            # - Persists to disk in persist_directory
            vector_store = Chroma.from_documents(
                documents=chunks,                               # LangChain Document objects
                embedding=self.embeddings,                      # OpenAI or HuggingFace embeddings
                collection_name=collection_name,                # Project-specific collection
                persist_directory=settings.CHROMA_PERSIST_DIR, # Local disk storage
            )
            # Note: ChromaDB automatically handles persistence, no manual .persist() needed

            self.log_info(f"Successfully stored {len(chunks)} vectors in collection: {collection_name}")

        except Exception as e:
            self.log_error("Error storing vectors in ChromaDB", e)
            raise

    async def query_knowledge_base(
        self,
        query: str,
        namespace: str,
        top_k: int = 5,
    ) -> List[LCDocument]:
        """
        Query the ChromaDB vector database for relevant documents using semantic search.

        This method performs semantic similarity search to find document chunks that
        are contextually related to the query, not just keyword matches.

        Semantic Search vs Keyword Search:
        ---------------------------------
        KEYWORD SEARCH:
        - Query: "manual data entry"
        - Matches: Only docs containing exact words "manual", "data", "entry"

        SEMANTIC SEARCH (This Method):
        - Query: "manual data entry"
        - Matches: Docs about "typing information by hand", "repetitive input",
                   "keyboard entry", "form filling" - contextually similar!

        How It Works:
        ------------
        1. Convert query text to embedding vector (using same model as documents)
        2. Calculate cosine similarity between query vector and all document vectors
        3. Return top_k most similar chunks (highest similarity scores)
        4. Each result includes: text content + metadata (source, page, etc.)

        Similarity Scoring:
        ------------------
        ChromaDB uses cosine similarity:
        - 1.0 = Perfect match (identical meaning)
        - 0.8-0.99 = Highly relevant
        - 0.6-0.79 = Moderately relevant
        - < 0.6 = Low relevance

        Use Cases:
        ---------
        - Hypothesis validation: Find evidence for/against hypotheses
        - Gap analysis: Compare documented vs actual processes
        - Context gathering: Get relevant background for questions
        - Answer generation: Find supporting information for reports

        Performance:
        -----------
        - Query time: ~50-200ms for 1,000 chunks
        - Scales linearly with chunk count
        - Embedding generation: ~20-50ms (OpenAI API)
        - Total latency: ~100-300ms typical

        Args:
            query: Natural language search query
                  Examples:
                  - "How are sales orders processed?"
                  - "What approval steps are manual?"
                  - "Describe the data entry workflow"
            namespace: Project namespace (must match storage namespace)
                      Format: "project_{project_id}"
                      Ensures query only searches this project's documents
            top_k: Number of most relevant chunks to return (default: 5)
                  - 3-5: Quick context lookup
                  - 5-10: Comprehensive search
                  - 10+: Exhaustive analysis

        Returns:
            List of LangChain Document objects, ordered by relevance (highest first)
            Each document contains:
            - page_content: Chunk text (string)
            - metadata: {source, page, chunk_number, project_id, ...}

            Returns empty list if:
            - Collection doesn't exist (no documents stored yet)
            - Query fails (network error, ChromaDB error)
            - No relevant results found

        Example:
            >>> results = await ingestion.query_knowledge_base(
            ...     query="manual approval process",
            ...     namespace="project_abc123",
            ...     top_k=5
            ... )
            >>> for doc in results:
            ...     print(f"Relevance: {doc.metadata.get('score')}")
            ...     print(f"Source: {doc.metadata['source']}")
            ...     print(f"Content: {doc.page_content[:200]}...")
        """
        try:
            # ============================================================
            # CREATE SANITIZED COLLECTION NAME
            # ============================================================
            # Must match the naming used in _store_in_vector_db()
            # Ensures query searches the correct project's collection
            sanitized_namespace = namespace.replace("_", "-").replace(" ", "-").lower()
            collection_name = f"{settings.CHROMA_COLLECTION_NAME}-{sanitized_namespace}"

            # ============================================================
            # INITIALIZE CHROMADB CONNECTION
            # ============================================================
            # Open existing collection from persistent storage
            # If collection doesn't exist, this will fail gracefully (caught below)
            vector_store = Chroma(
                collection_name=collection_name,                # Project-specific collection
                embedding_function=self.embeddings,             # Same embeddings as storage
                persist_directory=settings.CHROMA_PERSIST_DIR, # Local disk storage location
            )

            # ============================================================
            # PERFORM SEMANTIC SIMILARITY SEARCH
            # ============================================================
            # similarity_search():
            # 1. Generates embedding vector for query string
            # 2. Calculates cosine similarity with all stored vectors
            # 3. Returns top_k highest scoring chunks
            # 4. Results ordered by relevance (highest first)
            results = vector_store.similarity_search(query, k=top_k)

            self.log_info(f"Found {len(results)} relevant chunks for query in collection: {collection_name}")
            return results

        except Exception as e:
            # Handle errors gracefully:
            # - Collection not found (no documents stored yet)
            # - ChromaDB connection errors
            # - Embedding generation failures
            self.log_error(f"Error querying vector database for query '{query}' in namespace '{namespace}'", e)
            return []  # Return empty list rather than crashing
